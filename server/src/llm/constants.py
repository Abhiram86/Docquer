from io import BytesIO
from typing import List
import uuid
from pinecone import Index
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from PIL import Image
import easyocr
import numpy as np
from langchain_text_splitters import RecursiveCharacterTextSplitter
# from langchain_core.embeddings import FakeEmbeddings
# from langchain.embeddings import HuggingFaceEmbeddings
from sentence_transformers import SentenceTransformer
from pineconedb import pc, spec
from fastapi import HTTPException

embeddings_model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')


def normal_chat_main_content(name: str) -> str:
    return f"""
You are an assistant tasked with writing well-structured, high-quality blog-style content in Markdown format for {name}. 
Your job is to:
- Address the given prompt directly and concisely.
- Ensure proper Markdown syntax, especially for elements like tables, headers, bullet points, and code blocks.
- When using code blocks, ensure the correct language is specified (e.g., `python`, `javascript`, `bash`) to enable syntax highlighting.
- For tables, ensure they follow the correct syntax with `|` for columns and `-` for headers.
- Use a consistent style throughout the document, such as proper indentation, spacing after headers, and bullet point formatting.
- Make sure the content is clear, readable, and visually appealing.
- Avoid including any unnecessary information or explanations in the output.

Your output should be a valid, well-structured Markdown document with no unnecessary explanations or commentary. 
"""


def normal_chat_editor() -> str:
    return """
Your role is to inspect and refine the generated Markdown content. Follow these steps:

- **Tables**: Ensure all tables are correctly formatted with `|` for columns and `-` for headers. Check that the header row has at least three dashes (`---`) separating each column, and ensure proper alignment.
- **Code Blocks**: Inspect all code blocks to ensure they have the correct language name for syntax highlighting. For example, use `python` for Python code, `javascript` for JavaScript, etc. Ensure that only relevant code blocks are used (avoid using them for non-code sections).
- **Consistency**: Verify consistent indentation and spacing, especially after headers. Ensure that there is a blank line after each header and before any following content.
- **Formatting**: Ensure proper Markdown syntax, including the correct use of **bold** (`**text**`), **italic** (`*text*`), and **links** (`[text](url)`), when applicable.
- **Clarity**: Improve readability by ensuring clarity, structure, and flow. Eliminate any extraneous or redundant phrases.
- **No Preemble** : Avoid including any unnecessary information or explanations in the output.

Your output should only include the corrected and enhanced Markdown content, with no additional explanations.
"""


def getFileText(file: bytes, fileType: str) -> str:
    match fileType:
        case "text/plain":
            return readTXT(file)
        case "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return readDOCX(file)
        case "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return readPPTX(file)
        case "application/pdf":
            return readPDF(file)
        case "image/jpeg" | "image/png":
            return readImage(file)
        case _:
            return "nothing matched"


def readPDF(file: bytes) -> str:
    try:
        output = "This is a uploaded pdf file\n\n"
        pdf_stream = BytesIO(file)
        pdf_reader = PdfReader(pdf_stream)

        for i, page in enumerate(pdf_reader.pages):
            output += f"Page {i+1}:\n"
            # Extract text content
            page_text = page.extract_text()
            if page_text.strip():
                output += f"Uploaded PDF Content:\n{page_text}\n"

            # Extract and process images
            if '/XObject' in page['/Resources']:
                xObject = page['/Resources']['/XObject'].get_object()
                img_num = 1
                
                for obj in xObject:
                    if xObject[obj]['/Subtype'] == '/Image':
                        try:
                            image_data = xObject[obj].get_object()
                            # Convert image data to bytes
                            if image_data['/Filter'] == '/DCTDecode':
                                img_bytes = image_data._data
                                text_from_image = readImage(img_bytes)
                                if text_from_image and not text_from_image.isspace():
                                    output += f"\nImage {img_num} Text:\n{text_from_image}\n"
                                    img_num += 1
                        except Exception as img_err:
                            output += f"\nError processing image {img_num}: {str(img_err)}\n"
                            img_num += 1
            
            output += '-'*40 + "\n"

        return output
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing pdf file: {str(e)}")


def readPPTX(file: bytes) -> str:
    try:
        pptx_stream = BytesIO(file)
        presentation = Presentation(pptx_stream)
        output = "This is a uploaded pptx file\n\n"
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            output += f"Slide {slide_num}:\n"
            
            # Extract text content
            text_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_content.append(paragraph.text)
            
            if text_content:
                output += "Uploaded PPTX Paragraph:\n" + "\n".join(text_content) + "\n"
            
            # Process images
            img_num = 1
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        # Get image bytes
                        image_bytes = shape.image.blob
                        text_from_image = readImage(image_bytes)
                        if text_from_image and not text_from_image.isspace():
                            output += f"\nImage {img_num} Text:\n{text_from_image}\n"
                            img_num += 1
                    except Exception as img_err:
                        output += f"\nError processing image {img_num}: {str(img_err)}\n"
                        img_num += 1
            
            output += '-'*40 + "\n"
        
        return output
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing pptx file: {str(e)}")


def readTXT(file: bytes) -> str:
    return "This is a uploaded text file\n\n"+file.decode('utf-8')


def readDOCX(file: bytes) -> str:
    try:
        docx_stream = BytesIO(file)
        document = Document(docx_stream)
        output = "This is a uploaded docx file\n\n"
        
        # Track paragraph number for better organization
        para_num = 1
        img_num = 1
        
        # Iterate through the XML elements to preserve order
        for element in document.element.body:
            if element.tag.endswith("p"):  # Paragraph block (text)
                text = ''.join(node.text for node in element.iter() if node.text).strip()
                if text:
                    output += f"Uploaded DOCX Paragraph {para_num}: {text}\n"
                    para_num += 1
                    
            elif element.tag.endswith("drawing"):  # Image block
                try:
                    for rel in document.part.rels.values():
                        if "image" in rel.target_ref:  # Check for image reference
                            image_bytes = rel.target_part.blob
                            text_from_image = readImage(image_bytes)
                            if text_from_image and not text_from_image.isspace():
                                output += f"\nImage {img_num} Text:\n{text_from_image}\n"
                                img_num += 1
                except Exception as img_err:
                    output += f"\nError processing image {img_num}: {str(img_err)}\n"
                    img_num += 1
        
        return output.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing docx file: {str(e)}")


def readImage(file):
    image = Image.open(BytesIO(file))
    reader = easyocr.Reader(['en'])
    array = np.array(image)
    res = reader.readtext(array)
    text = "This is a uploaded image file\n\n" + "".join([r[1] for r in res])
    return text


def split_into_chunks(data: str) -> List[str]:
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        length_function=len,
        is_separator_regex=False,
    )
    texts = text_splitter.split_text(data)
    return texts


def get_index(conv_id: str) -> Index:
    index_name = f"docquer-{conv_id}"
    try:
        index = pc.Index(index_name)
        # Try to describe index to ensure it exists and is ready
        index.describe_index_stats()
        return index
    except Exception as e:
        print(f"Error getting index: {e}")
        raise HTTPException(status_code=500, detail="Index not ready or doesn't exist")


def insert_data(conv_id: str, data: List[str], replace: bool = False):
    if not data:
        return
        
    try:
        index = get_index(conv_id)
        
        # If replace is True, delete existing vectors
        if replace:
            try:
                index.delete(delete_all=True)
            except Exception as e:
                print(f"Error deleting vectors: {e}")
                
        embeddings = embeddings_model.encode(data)
        vectors = [{
            "id": str(uuid.uuid4()),
            "values": embedding.tolist(),
            "metadata": {"text": data[i]}
        } for i, embedding in enumerate(embeddings)]
        
        # Split vectors into smaller batches
        batch_size = 100
        for i in range(0, len(vectors), batch_size):
            batch = vectors[i:i + batch_size]
            index.upsert(vectors=batch)
            
    except Exception as e:
        print(f"Error inserting data: {e}")
        raise HTTPException(status_code=500, detail=f"Error inserting data: {str(e)}")


def init_vector_db(chunks: List[str], conv_id: str) -> Index:
    if not chunks:
        raise HTTPException(status_code=400, detail="No data to index")

    index_name = f"docquer-{conv_id}"
    try:
        # Check if index exists and is ready
        try:
            index = pc.Index(index_name)
            index.describe_index_stats()
        except:
            # Create new index if it doesn't exist or isn't ready
            if index_name in pc.list_indexes():
                pc.delete_index(index_name)
            pc.create_index(
                name=index_name,
                dimension=384,
                metric='cosine',
                spec=spec
            )
            # Wait for index to be ready
            index = pc.Index(index_name)
            retries = 0
            while retries < 5:
                try:
                    index.describe_index_stats()
                    break
                except:
                    retries += 1
                    import time
                    time.sleep(2)

        # Insert data in batches
        insert_data(conv_id, chunks)
        return index
    except Exception as e:
        print(f"Error initializing vector DB: {e}")
        raise HTTPException(status_code=500, detail=f"Error initializing vector DB: {str(e)}")
